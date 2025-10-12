"""
RAG Client for document processing, vector database operations, and knowledge graph management.

This module provides comprehensive functionality for:
- Document processing (text extraction, chunking)
- Vector database operations (ChromaDB integration)
- Knowledge graph management (Neo4j integration)
- Semantic search and retrieval
- Entity extraction and relationship discovery
"""

import os
import hashlib
import mimetypes
import logging
from typing import List, Dict, Any, Optional, Tuple, Set
from pathlib import Path
from datetime import datetime, timezone
from collections import defaultdict
import re

# Document processing
import fitz  # PyMuPDF for PDF processing
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation

# Vector database and embeddings
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer

# Graph database
from neo4j import GraphDatabase
import spacy

# Text processing
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document

logger = logging.getLogger(__name__)

class RAGClient:
    """Client for RAG operations including document processing and vector search."""

    def __init__(self,
                 chroma_persist_directory: str = "./chroma_db",
                 embedding_model: str = "all-MiniLM-L6-v2",
                 neo4j_uri: Optional[str] = None,
                 neo4j_user: Optional[str] = None,
                 neo4j_password: Optional[str] = None):
        """Initialize RAG client with database connections.

        Args:
            chroma_persist_directory: Directory for ChromaDB persistence
            embedding_model: SentenceTransformer model name
            neo4j_uri: Neo4j database URI
            neo4j_user: Neo4j username
            neo4j_password: Neo4j password
        """
        self.persist_directory = chroma_persist_directory
        self.embedding_model_name = embedding_model

        # Initialize ChromaDB
        self.chroma_client = chromadb.PersistentClient(
            path=chroma_persist_directory,
            settings=Settings(anonymized_telemetry=False)
        )

        # Initialize embedding model
        self.embedding_model = SentenceTransformer(embedding_model)

        # Initialize text splitter for document chunking
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )

        # Initialize Neo4j connection if credentials provided
        self.neo4j_driver = None
        if all([neo4j_uri, neo4j_user, neo4j_password]):
            try:
                self.neo4j_driver = GraphDatabase.driver(
                    neo4j_uri, auth=(neo4j_user, neo4j_password)
                )
                logger.info("Connected to Neo4j database")
            except Exception as e:
                logger.warning(f"Failed to connect to Neo4j: {e}")

        # Initialize spaCy for NER (download model if needed)
        try:
            self.nlp = spacy.load("en_core_web_sm")
        except OSError:
            logger.warning("spaCy model 'en_core_web_sm' not found. Install with: python -m spacy download en_core_web_sm")
            self.nlp = None

        # File type handlers
        self.file_handlers = {
            '.txt': self._extract_text_from_txt,
            '.md': self._extract_text_from_txt,
            '.py': self._extract_text_from_txt,
            '.js': self._extract_text_from_txt,
            '.json': self._extract_text_from_txt,
            '.yaml': self._extract_text_from_txt,
            '.yml': self._extract_text_from_txt,
            '.pdf': self._extract_text_from_pdf,
            '.docx': self._extract_text_from_docx,
            '.xlsx': self._extract_text_from_xlsx,
            '.pptx': self._extract_text_from_pptx,
        }

        # Supported file extensions
        self.supported_extensions = set(self.file_handlers.keys())

    def scan_folder(self, folder_path: str,
                   extensions: Optional[List[str]] = None,
                   exclude_patterns: Optional[List[str]] = None,
                   max_file_size_mb: int = 50) -> List[Dict[str, Any]]:
        """Recursively scan folder for supported documents.

        Args:
            folder_path: Path to folder to scan
            extensions: List of file extensions to include (default: all supported)
            exclude_patterns: List of regex patterns to exclude files/folders
            max_file_size_mb: Maximum file size in MB

        Returns:
            List of file information dictionaries
        """
        if extensions is None:
            extensions = list(self.supported_extensions)

        exclude_patterns = exclude_patterns or []
        compiled_patterns = [re.compile(pattern) for pattern in exclude_patterns]

        files = []
        folder_path = Path(folder_path)

        if not folder_path.exists():
            raise FileNotFoundError(f"Folder not found: {folder_path}")

        for file_path in folder_path.rglob("*"):
            if not file_path.is_file():
                continue

            # Check file extension
            if file_path.suffix.lower() not in extensions:
                continue

            # Check exclude patterns
            if any(pattern.search(str(file_path)) for pattern in compiled_patterns):
                continue

            # Check file size
            try:
                file_size = file_path.stat().st_size
                if file_size > max_file_size_mb * 1024 * 1024:
                    logger.warning(f"Skipping large file: {file_path} ({file_size / 1024 / 1024:.1f}MB)")
                    continue
            except OSError:
                continue

            # Get file info
            try:
                stat = file_path.stat()
                files.append({
                    'path': str(file_path),
                    'name': file_path.name,
                    'extension': file_path.suffix.lower(),
                    'size_bytes': stat.st_size,
                    'modified_time': datetime.fromtimestamp(stat.st_mtime, timezone.utc),
                    'relative_path': str(file_path.relative_to(folder_path))
                })
            except OSError as e:
                logger.warning(f"Error reading file info for {file_path}: {e}")
                continue

        logger.info(f"Found {len(files)} files in {folder_path}")
        return files

    def extract_text(self, file_path: str) -> Dict[str, Any]:
        """Extract text content from a file.

        Args:
            file_path: Path to the file

        Returns:
            Dictionary with extracted text and metadata
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        extension = file_path.suffix.lower()
        if extension not in self.file_handlers:
            raise ValueError(f"Unsupported file type: {extension}")

        try:
            # Extract text using appropriate handler
            text = self.file_handlers[extension](file_path)

            # Calculate file hash for deduplication
            file_hash = self._calculate_file_hash(file_path)

            # Get file metadata
            stat = file_path.stat()

            return {
                'file_path': str(file_path),
                'file_name': file_path.name,
                'file_extension': extension,
                'file_size': stat.st_size,
                'modified_time': datetime.fromtimestamp(stat.st_mtime, timezone.utc),
                'file_hash': file_hash,
                'text_content': text,
                'text_length': len(text),
                'extraction_time': datetime.now(timezone.utc)
            }

        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            raise

    def chunk_documents(self, documents: List[Dict[str, Any]],
                       chunk_size: int = 1000,
                       chunk_overlap: int = 200) -> List[Dict[str, Any]]:
        """Split documents into chunks for embedding.

        Args:
            documents: List of document dictionaries with text content
            chunk_size: Target size for each chunk
            chunk_overlap: Overlap between chunks

        Returns:
            List of chunk dictionaries
        """
        # Update text splitter with new parameters
        if chunk_size != self.text_splitter._chunk_size or chunk_overlap != self.text_splitter._chunk_overlap:
            self.text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                length_function=len,
                separators=["\n\n", "\n", " ", ""]
            )

        chunks = []

        for doc in documents:
            text_content = doc.get('text_content', '')
            if not text_content.strip():
                continue

            # Split text into chunks
            text_chunks = self.text_splitter.split_text(text_content)

            # Create chunk objects
            for i, chunk_text in enumerate(text_chunks):
                chunk_id = f"{doc.get('file_hash', 'unknown')}_{i}"

                chunks.append({
                    'chunk_id': chunk_id,
                    'chunk_index': i,
                    'chunk_text': chunk_text,
                    'chunk_length': len(chunk_text),
                    'source_file': doc.get('file_path', ''),
                    'source_file_name': doc.get('file_name', ''),
                    'source_extension': doc.get('file_extension', ''),
                    'source_hash': doc.get('file_hash', ''),
                    'source_modified_time': doc.get('modified_time'),
                    'chunk_created_time': datetime.now(timezone.utc)
                })

        logger.info(f"Created {len(chunks)} chunks from {len(documents)} documents")
        return chunks

    def generate_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Generate embeddings for text chunks.

        Args:
            texts: List of text strings to embed

        Returns:
            List of embedding vectors
        """
        if not texts:
            return []

        try:
            embeddings = self.embedding_model.encode(
                texts,
                batch_size=32,
                show_progress_bar=len(texts) > 100
            )
            return embeddings.tolist()
        except Exception as e:
            logger.error(f"Error generating embeddings: {e}")
            raise

    def detect_duplicates(self, documents: List[Dict[str, Any]],
                         similarity_threshold: float = 0.95) -> List[List[int]]:
        """Detect duplicate documents based on content similarity.

        Args:
            documents: List of document dictionaries
            similarity_threshold: Similarity threshold for duplicate detection

        Returns:
            List of duplicate groups (each group is a list of document indices)
        """
        if len(documents) < 2:
            return []

        # Extract text content for comparison
        texts = [doc.get('text_content', '') for doc in documents]

        # Generate embeddings for duplicate detection
        embeddings = self.generate_embeddings(texts)

        # Find duplicates using cosine similarity
        duplicates = []
        processed = set()

        for i in range(len(embeddings)):
            if i in processed:
                continue

            group = [i]
            processed.add(i)

            for j in range(i + 1, len(embeddings)):
                if j in processed:
                    continue

                # Calculate cosine similarity
                similarity = self._cosine_similarity(embeddings[i], embeddings[j])

                if similarity >= similarity_threshold:
                    group.append(j)
                    processed.add(j)

            if len(group) > 1:
                duplicates.append(group)

        logger.info(f"Found {len(duplicates)} duplicate groups")
        return duplicates

    def build_vector_index(self, collection_name: str,
                          chunks: List[Dict[str, Any]],
                          metadata_fields: Optional[List[str]] = None) -> str:
        """Build vector index from document chunks.

        Args:
            collection_name: Name for the vector collection
            chunks: List of chunk dictionaries
            metadata_fields: List of metadata fields to include

        Returns:
            Collection name
        """
        if not chunks:
            raise ValueError("No chunks provided for indexing")

        # Get or create collection
        try:
            collection = self.chroma_client.get_collection(collection_name)
            logger.info(f"Using existing collection: {collection_name}")
        except:
            collection = self.chroma_client.create_collection(
                name=collection_name,
                metadata={"description": f"Vector index created {datetime.now(timezone.utc)}"}
            )
            logger.info(f"Created new collection: {collection_name}")

        # Prepare data for indexing
        texts = [chunk['chunk_text'] for chunk in chunks]
        ids = [chunk['chunk_id'] for chunk in chunks]

        # Generate embeddings
        embeddings = self.generate_embeddings(texts)

        # Prepare metadata
        metadatas = []
        default_fields = ['source_file', 'source_file_name', 'chunk_index', 'chunk_length']
        fields_to_include = metadata_fields or default_fields

        for chunk in chunks:
            metadata = {}
            for field in fields_to_include:
                if field in chunk:
                    value = chunk[field]
                    # Convert datetime to string for storage
                    if isinstance(value, datetime):
                        value = value.isoformat()
                    metadata[field] = value
            metadatas.append(metadata)

        # Add documents to collection
        collection.add(
            documents=texts,
            embeddings=embeddings,
            metadatas=metadatas,
            ids=ids
        )

        logger.info(f"Added {len(chunks)} chunks to collection {collection_name}")
        return collection_name

    def semantic_search(self, collection_name: str,
                       query: str,
                       n_results: int = 10,
                       filter_metadata: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
        """Perform semantic search on vector collection.

        Args:
            collection_name: Name of the collection to search
            query: Search query text
            n_results: Number of results to return
            filter_metadata: Optional metadata filters

        Returns:
            List of search results with scores and metadata
        """
        try:
            collection = self.chroma_client.get_collection(collection_name)
        except:
            raise ValueError(f"Collection not found: {collection_name}")

        # Generate query embedding
        query_embedding = self.generate_embeddings([query])[0]

        # Perform search
        results = collection.query(
            query_embeddings=[query_embedding],
            n_results=n_results,
            where=filter_metadata
        )

        # Format results
        search_results = []
        for i in range(len(results['ids'][0])):
            search_results.append({
                'id': results['ids'][0][i],
                'text': results['documents'][0][i],
                'score': float(results['distances'][0][i]),
                'metadata': results['metadatas'][0][i]
            })

        return search_results

    def update_index(self, collection_name: str,
                    chunks: List[Dict[str, Any]]) -> int:
        """Update vector index with new chunks.

        Args:
            collection_name: Name of the collection to update
            chunks: List of new chunk dictionaries

        Returns:
            Number of chunks added
        """
        if not chunks:
            return 0

        try:
            collection = self.chroma_client.get_collection(collection_name)
        except:
            logger.info(f"Collection {collection_name} not found, creating new one")
            return len(self.build_vector_index(collection_name, chunks))

        # Check for existing chunks to avoid duplicates
        existing_ids = set()
        try:
            # Get existing IDs (this might fail if collection is empty)
            peek_result = collection.peek()
            if peek_result and peek_result.get('ids'):
                existing_ids = set(peek_result['ids'])
        except:
            pass

        # Filter out existing chunks
        new_chunks = [chunk for chunk in chunks if chunk['chunk_id'] not in existing_ids]

        if not new_chunks:
            logger.info("No new chunks to add")
            return 0

        # Add new chunks using build_vector_index logic
        texts = [chunk['chunk_text'] for chunk in new_chunks]
        ids = [chunk['chunk_id'] for chunk in new_chunks]
        embeddings = self.generate_embeddings(texts)

        # Prepare metadata
        metadatas = []
        for chunk in new_chunks:
            metadata = {
                'source_file': chunk.get('source_file', ''),
                'source_file_name': chunk.get('source_file_name', ''),
                'chunk_index': chunk.get('chunk_index', 0),
                'chunk_length': chunk.get('chunk_length', 0)
            }
            metadatas.append(metadata)

        collection.add(
            documents=texts,
            embeddings=embeddings,
            metadatas=metadatas,
            ids=ids
        )

        logger.info(f"Added {len(new_chunks)} new chunks to collection {collection_name}")
        return len(new_chunks)

    def manage_collections(self) -> Dict[str, Any]:
        """Get information about all vector collections.

        Returns:
            Dictionary with collection information
        """
        collections = self.chroma_client.list_collections()

        collection_info = {}
        for collection in collections:
            try:
                count = collection.count()
                collection_info[collection.name] = {
                    'name': collection.name,
                    'count': count,
                    'metadata': collection.metadata
                }
            except Exception as e:
                logger.warning(f"Error getting info for collection {collection.name}: {e}")
                collection_info[collection.name] = {
                    'name': collection.name,
                    'count': 0,
                    'error': str(e)
                }

        return collection_info

    def extract_entities(self, text: str) -> List[Dict[str, Any]]:
        """Extract named entities from text using spaCy NER.

        Args:
            text: Text to process

        Returns:
            List of entity dictionaries
        """
        if not self.nlp:
            raise RuntimeError("spaCy model not available. Install with: python -m spacy download en_core_web_sm")

        doc = self.nlp(text)
        entities = []

        for ent in doc.ents:
            entities.append({
                'text': ent.text,
                'label': ent.label_,
                'description': spacy.explain(ent.label_),
                'start_char': ent.start_char,
                'end_char': ent.end_char,
                'confidence': getattr(ent, 'score', 1.0)
            })

        return entities

    def build_knowledge_graph(self, documents: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Build knowledge graph from documents using Neo4j.

        Args:
            documents: List of document dictionaries

        Returns:
            Dictionary with graph statistics
        """
        if not self.neo4j_driver:
            raise RuntimeError("Neo4j connection not available")

        entities_by_doc = {}
        relationships = []

        # Extract entities from each document
        for doc in documents:
            doc_id = doc.get('file_hash', doc.get('file_path', ''))
            text = doc.get('text_content', '')

            if not text:
                continue

            entities = self.extract_entities(text)
            entities_by_doc[doc_id] = entities

        # Create graph nodes and relationships
        with self.neo4j_driver.session() as session:
            # Clear existing graph (optional)
            # session.run("MATCH (n) DETACH DELETE n")

            nodes_created = 0
            relationships_created = 0

            # Create document nodes
            for doc in documents:
                doc_id = doc.get('file_hash', doc.get('file_path', ''))
                session.run("""
                    MERGE (d:Document {id: $doc_id})
                    SET d.name = $name,
                        d.path = $path,
                        d.size = $size,
                        d.modified = $modified
                """, {
                    'doc_id': doc_id,
                    'name': doc.get('file_name', ''),
                    'path': doc.get('file_path', ''),
                    'size': doc.get('file_size', 0),
                    'modified': doc.get('modified_time', '').isoformat() if doc.get('modified_time') else ''
                })
                nodes_created += 1

            # Create entity nodes and relationships
            for doc_id, entities in entities_by_doc.items():
                for entity in entities:
                    # Create entity node
                    session.run("""
                        MERGE (e:Entity {text: $text, label: $label})
                        SET e.description = $description
                    """, {
                        'text': entity['text'],
                        'label': entity['label'],
                        'description': entity['description']
                    })

                    # Create relationship between document and entity
                    session.run("""
                        MATCH (d:Document {id: $doc_id})
                        MATCH (e:Entity {text: $text, label: $label})
                        MERGE (d)-[r:CONTAINS]->(e)
                        SET r.start_char = $start_char,
                            r.end_char = $end_char,
                            r.confidence = $confidence
                    """, {
                        'doc_id': doc_id,
                        'text': entity['text'],
                        'label': entity['label'],
                        'start_char': entity['start_char'],
                        'end_char': entity['end_char'],
                        'confidence': entity['confidence']
                    })
                    relationships_created += 1

        return {
            'documents_processed': len(documents),
            'total_entities': sum(len(entities) for entities in entities_by_doc.values()),
            'nodes_created': nodes_created,
            'relationships_created': relationships_created
        }

    def find_relationships(self, entity1: str, entity2: str) -> List[Dict[str, Any]]:
        """Find relationships between entities in the knowledge graph.

        Args:
            entity1: First entity text
            entity2: Second entity text

        Returns:
            List of relationship paths
        """
        if not self.neo4j_driver:
            raise RuntimeError("Neo4j connection not available")

        with self.neo4j_driver.session() as session:
            # Find direct relationships
            result = session.run("""
                MATCH (e1:Entity {text: $entity1})-[r1:CONTAINS]-(d:Document)-[r2:CONTAINS]-(e2:Entity {text: $entity2})
                RETURN d.name as document, d.path as path,
                       e1.text as entity1, e1.label as label1,
                       e2.text as entity2, e2.label as label2,
                       r1.confidence as conf1, r2.confidence as conf2
                ORDER BY (r1.confidence + r2.confidence) DESC
                LIMIT 10
            """, {'entity1': entity1, 'entity2': entity2})

            relationships = []
            for record in result:
                relationships.append({
                    'document': record['document'],
                    'document_path': record['path'],
                    'entity1': record['entity1'],
                    'entity1_label': record['label1'],
                    'entity2': record['entity2'],
                    'entity2_label': record['label2'],
                    'confidence_score': (record['conf1'] + record['conf2']) / 2,
                    'relationship_type': 'co_occurrence'
                })

        return relationships

    def graph_query(self, cypher_query: str) -> List[Dict[str, Any]]:
        """Execute custom Cypher query on knowledge graph.

        Args:
            cypher_query: Cypher query string

        Returns:
            List of query results
        """
        if not self.neo4j_driver:
            raise RuntimeError("Neo4j connection not available")

        with self.neo4j_driver.session() as session:
            result = session.run(cypher_query)
            return [record.data() for record in result]

    # Private helper methods

    def _extract_text_from_txt(self, file_path: Path) -> str:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1', errors='ignore') as f:
                return f.read()

    def _extract_text_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file."""
        text = ""
        try:
            with fitz.open(str(file_path)) as pdf_document:
                for page_num in range(pdf_document.page_count):
                    page = pdf_document[page_num]
                    text += page.get_text()
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {e}")
            return ""

    def _extract_text_from_docx(self, file_path: Path) -> str:
        """Extract text from DOCX file."""
        try:
            doc = DocxDocument(str(file_path))
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {e}")
            return ""

    def _extract_text_from_xlsx(self, file_path: Path) -> str:
        """Extract text from XLSX file."""
        try:
            workbook = openpyxl.load_workbook(str(file_path), data_only=True)
            text = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text.append(f"Sheet: {sheet_name}")

                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text.append(row_text)

            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting text from XLSX {file_path}: {e}")
            return ""

    def _extract_text_from_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX file."""
        try:
            prs = Presentation(str(file_path))
            text = []

            for i, slide in enumerate(prs.slides):
                text.append(f"Slide {i + 1}:")

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)

            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting text from PPTX {file_path}: {e}")
            return ""

    def _calculate_file_hash(self, file_path: Path) -> str:
        """Calculate SHA-256 hash of file for deduplication."""
        hash_sha256 = hashlib.sha256()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_sha256.update(chunk)
            return hash_sha256.hexdigest()
        except Exception as e:
            logger.error(f"Error calculating hash for {file_path}: {e}")
            return str(hash(str(file_path)))

    def _cosine_similarity(self, vec1: List[float], vec2: List[float]) -> float:
        """Calculate cosine similarity between two vectors."""
        import math

        dot_product = sum(a * b for a, b in zip(vec1, vec2))
        magnitude1 = math.sqrt(sum(a * a for a in vec1))
        magnitude2 = math.sqrt(sum(b * b for b in vec2))

        if magnitude1 == 0 or magnitude2 == 0:
            return 0.0

        return dot_product / (magnitude1 * magnitude2)

    def __del__(self):
        """Clean up database connections."""
        if self.neo4j_driver:
            try:
                self.neo4j_driver.close()
            except:
                pass